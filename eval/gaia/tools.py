"""Real, read-only tools for the official GAIA agent loop."""

from __future__ import annotations

import ast
import base64
import csv
import json
import math
import operator
import os
from dataclasses import dataclass
from pathlib import Path
from typing import Any


@dataclass
class ToolResult:
    text: str
    error: str | None = None
    image_media_type: str | None = None
    image_base64: str | None = None

    def content_blocks(self) -> list[dict[str, Any]]:
        blocks: list[dict[str, Any]] = [{"type": "text", "text": self.text}]
        if self.image_media_type and self.image_base64:
            blocks.append(
                {
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": self.image_media_type,
                        "data": self.image_base64,
                    },
                }
            )
        return blocks


_BINOPS = {
    ast.Add: operator.add,
    ast.Sub: operator.sub,
    ast.Mult: operator.mul,
    ast.Div: operator.truediv,
    ast.FloorDiv: operator.floordiv,
    ast.Mod: operator.mod,
    ast.Pow: operator.pow,
}
_UNARY = {ast.UAdd: operator.pos, ast.USub: operator.neg}
_CONSTANTS = {"pi": math.pi, "e": math.e}


def calculate(expression: str) -> ToolResult:
    """Evaluate arithmetic through a restricted AST, never eval/exec."""

    def visit(node: ast.AST) -> float | int:
        if isinstance(node, ast.Expression):
            return visit(node.body)
        if isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
            return node.value
        if isinstance(node, ast.Name) and node.id in _CONSTANTS:
            return _CONSTANTS[node.id]
        if isinstance(node, ast.BinOp) and type(node.op) in _BINOPS:
            left, right = visit(node.left), visit(node.right)
            if isinstance(node.op, ast.Pow) and abs(right) > 100:
                raise ValueError("exponent outside safe range")
            return _BINOPS[type(node.op)](left, right)
        if isinstance(node, ast.UnaryOp) and type(node.op) in _UNARY:
            return _UNARY[type(node.op)](visit(node.operand))
        raise ValueError(f"unsupported calculator syntax: {type(node).__name__}")

    try:
        parsed = ast.parse(expression.strip(), mode="eval")
        value = visit(parsed)
        if isinstance(value, float) and not math.isfinite(value):
            raise ValueError("non-finite result")
        return ToolResult(text=str(value))
    except Exception as exc:
        return ToolResult(text=f"Calculator error: {exc}", error="calculator_error")


def web_search(client: Any, query: str, max_results: int = 5) -> ToolResult:
    """Execute a real Tavily search and return bounded source excerpts."""
    try:
        payload = client.search(
            query=query,
            search_depth="advanced",
            max_results=max(1, min(int(max_results), 8)),
            include_answer=True,
        )
        rows = []
        if payload.get("answer"):
            rows.append(f"Answer summary: {payload['answer']}")
        for index, item in enumerate(payload.get("results", []), 1):
            rows.append(
                f"[{index}] {item.get('title', 'Untitled')}\n"
                f"URL: {item.get('url', '')}\n"
                f"Excerpt: {item.get('content', '')[:4000]}"
            )
        if not rows:
            return ToolResult(text="No search results returned.", error="retrieval_error")
        return ToolResult(text="\n\n".join(rows))
    except Exception as exc:
        return ToolResult(text=f"Search error: {exc}", error="search_error")


def transcribe_audio(path: Path, client: Any | None = None) -> ToolResult:
    """Transcribe a GAIA audio attachment through the approved OpenAI boundary."""
    if client is None:
        if not os.getenv("OPENAI_API_KEY"):
            return ToolResult(
                "OPENAI_API_KEY is required for audio transcription.",
                "file_read_error",
            )
        from openai import OpenAI

        client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
    try:
        with path.open("rb") as handle:
            response = client.audio.transcriptions.create(
                model="gpt-4o-mini-transcribe",
                file=handle,
            )
        text = response.text.strip()
        return ToolResult(text=text or "[Audio transcription was empty]")
    except Exception as exc:
        return ToolResult(f"Audio transcription error: {exc}", "file_read_error")


def read_file(path_value: str, attachment_root: Path, max_chars: int = 50_000) -> ToolResult:
    """Read a GAIA attachment without permitting access outside its cache root."""
    path = Path(path_value)
    if path.is_absolute():
        try:
            path.relative_to(attachment_root)
        except ValueError:
            return ToolResult("File access denied: outside GAIA attachment root.", "file_access_error")
    else:
        if ".." in path.parts:
            return ToolResult("File access denied: parent traversal.", "file_access_error")
        path = attachment_root / path
    # Hugging Face snapshots intentionally symlink into their trusted blob cache.
    # Lexical confinement above allows those links without allowing user traversal.
    if not path.is_file():
        return ToolResult(f"File not found: {path.name}", "file_not_found")

    suffix = path.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".py", ".json", ".xml", ".html", ".htm"}:
            text = path.read_text(errors="replace")
        elif suffix in {".csv", ".tsv"}:
            dialect = "excel-tab" if suffix == ".tsv" else "excel"
            with path.open(newline="", errors="replace") as handle:
                rows = list(csv.reader(handle, dialect=dialect))
            text = "\n".join("\t".join(cell for cell in row) for row in rows)
        elif suffix == ".pdf":
            from pypdf import PdfReader

            reader = PdfReader(path)
            text = "\n\n".join(
                f"--- Page {index} ---\n{page.extract_text() or ''}"
                for index, page in enumerate(reader.pages, 1)
            )
        elif suffix in {".xlsx", ".xlsm"}:
            from openpyxl import load_workbook

            workbook = load_workbook(path, read_only=True, data_only=True)
            chunks = []
            for sheet in workbook.worksheets:
                chunks.append(f"--- Sheet: {sheet.title} ---")
                for row in sheet.iter_rows(values_only=True):
                    chunks.append("\t".join("" if value is None else str(value) for value in row))
            text = "\n".join(chunks)
        elif suffix == ".docx":
            from docx import Document

            document = Document(path)
            text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        elif suffix == ".pptx":
            from pptx import Presentation

            presentation = Presentation(path)
            chunks = []
            for index, slide in enumerate(presentation.slides, 1):
                chunks.append(f"--- Slide {index} ---")
                chunks.extend(
                    shape.text
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                )
            text = "\n".join(chunks)
        elif suffix in {".mp3", ".wav", ".m4a", ".flac", ".ogg"}:
            return transcribe_audio(path)
        elif suffix in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
            media = {
                ".png": "image/png",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".webp": "image/webp",
                ".gif": "image/gif",
            }[suffix]
            encoded = base64.b64encode(path.read_bytes()).decode("ascii")
            return ToolResult(
                text=f"Image attachment loaded: {path.name}",
                image_media_type=media,
                image_base64=encoded,
            )
        else:
            return ToolResult(
                f"Unsupported attachment format: {suffix or '[none]'}",
                "unsupported_file_type",
            )
        truncated = len(text) > max_chars
        text = text[:max_chars]
        if truncated:
            text += f"\n[Truncated at {max_chars} characters]"
        return ToolResult(text=text or "[File contained no extractable text]")
    except Exception as exc:
        return ToolResult(f"File-read error for {path.name}: {exc}", "file_read_error")
